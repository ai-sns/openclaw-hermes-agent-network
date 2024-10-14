from selenium import webdriver
from selenium.webdriver.chrome.service import Service as ChromeService
from webdriver_manager.chrome import ChromeDriverManager
from selenium.webdriver.common.keys import Keys
from selenium.webdriver.common.by import By
import pandas as pd
from pptx import Presentation
from pptx.util import Inches
import os
from pathlib import Path
import shutil
import datetime




class getDatas():

    def __init__(self, company,DRIVER):
        self.company = company
        if DRIVER == None:
            self.DRIVER =  webdriver.Chrome(service=ChromeService(ChromeDriverManager().install()))
        else:
            self.DRIVER = DRIVER

    def webScraping(self):

        DRIVER =self.DRIVER
        DRIVER.get('https://www.google.com')
        DRIVER.find_element('xpath','//*[@id="APjFqb"]').send_keys(f'{self.company} stocks price', Keys.ENTER)
        name = DRIVER.find_element(By.CLASS_NAME,'aMEhee').text
        value = DRIVER.find_element('xpath','//*[@id="knowledge-finance-wholepage__entity-summary"]/div[3]/g-card-section/div/g-card-section/div[2]/div[1]/span[1]/span/span[1]').text
        coin = DRIVER.find_element('xpath','//*[@id="knowledge-finance-wholepage__entity-summary"]/div[3]/g-card-section/div/g-card-section/div[2]/div[1]/span[1]/span/span[2]').text

        return name, value, coin

class sheets:
    def __init__(self, name, value, coin):
        self.name = name
        self.value = value
        self.coin = coin

    def importSheet(self):
        try:
            # Insert your andress of directory
            directory = os.path.join(Path(__file__).resolve().parent, "temp", "market")

            NOW = datetime.datetime.now()
            CSV = f'{directory}-{NOW.month}-{NOW.year}.csv'

            df = pd.DataFrame({'Company': self.name, 'Value': self.value, 'Coin': self.coin})
            df.to_csv(CSV, index=False)
        except Exception as e:  # 捕获异常并将其赋值给变量e
            print(f"An error occurred: {e}")  # 打印异常信息

        return df

class slide:
    def __init__(self,dataFrame):
        self.dataFrame = dataFrame

    def importSlide(self):
        df = pd.read_csv(self.dataFrame)

        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        title = slide.shapes.title
        rows = len(df.index)+1
        cols = len(df.columns)
        x = Inches(2.5)
        y = Inches(2.0)
        width = Inches(5.0)
        height = Inches(1.0)
        table = slide.shapes.add_table(rows,cols,x,y,width,height).table
        title.text = "Summary of Market"

        for company in range(cols):
            table.cell(0,company).text=str(df.columns[company])

        for company in range(len(df.index)):
            for values in range(cols):
                table.cell(company+1,values).text=str(df.iloc[company,values])
        # Insert your andress of directory
        directory = os.path.join(Path(__file__).resolve().parent, "temp", "market")

        NOW = datetime.datetime.now()
        PPTX = f'{directory}-{NOW.month}-{NOW.year}.pptx'
        prs.save(PPTX)

class backup:
    def __init__(self, csv, pptx):
        self.csv = csv
        self.pptx = pptx

    def save(self):
        if (os.path.exists(self.csv)):
            directory = os.path.join(Path(__file__).resolve().parent, "temp", "market")
            shutil.move(self.csv,f'{directory[2:]}Backup')
            shutil.move(self.pptx,f'{directory[2:]}Backup')

if __name__== '__main__':
    pass
