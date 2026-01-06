import os
from typing import Union, List
import jieba
# from langchain_text_splitters import RecursiveCharacterTextSplitter, CharacterTextSplitter

os.environ['HUGGINGFACEHUB_API_TOKEN'] = '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'
os.environ[
    "OPENAI_API_KEY"] = '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'
os.environ["GOOGLE_CSE_ID"] = "53b9c3fd76d8d4cbb"
os.environ["GOOGLE_API_KEY"] = "AIzaSyAYEpRPu24tU41bFn4QQB_2cZFmlOZxEEE"
# from langchain_community.document_loaders import TextLoader, powerpoint, word_document, excel, PyPDFLoader, \
#     markdown, html, csv_loader
from langchain_openai import OpenAIEmbeddings
# from langchain.text_splitter import TokenTextSplitter
from langchain_community.vectorstores import Chroma
from langchain_community.embeddings import HuggingFaceEmbeddings
from util import image_to_base64


"""
 pip install chromadb==0.4.24
 pip install tiktoken
"""


def savevector(filepath, persist_directory, embedding_model_name, emb_type="openai", chunk_size=500, chunk_overlap=20):
    # 指定chroma持久化的目录，当我们不知道目录时,chroma会将数据存储在内存中，随着程序的关闭就会删除
    # persist_directory = "C:\\dev\\ai-sns\\PyTalk\\pytalk\\vector_store"
    # 按目录加载文档
    # loader = DirectoryLoader('C:\\0资料\\12.Omniverse\\青田项目\\经商局\\km\\cleaned\\', glob='**/*.txt')
    # docs = loader.load()
    # 加载单个文档 可以自由选择
    # loader = TextLoader(filepath, encoding='utf8')
    ext = os.path.splitext(filepath)[1].lower()  # 获取文件扩展名并转为小写
    # loaders = {
    #     '.txt': lambda path: TextLoader(path, encoding='utf8'),
    #     '.js': lambda path: TextLoader(path, encoding='utf8'),
    #     '.sql': lambda path: TextLoader(path, encoding='utf8'),
    #     '.pdf': PyPDFLoader,
    #     '.docx': word_document.UnstructuredWordDocumentLoader,
    #     '.xls': excel.UnstructuredExcelLoader,
    #     '.xlsx': excel.UnstructuredExcelLoader,
    #     '.csv': csv_loader.UnstructuredCSVLoader,
    #     '.pptx': powerpoint.UnstructuredPowerPointLoader,
    #     '.md': markdown.UnstructuredMarkdownLoader,
    #     '.html': html.UnstructuredHTMLLoader,
    #     '.htm': html.UnstructuredHTMLLoader,
    # }
    #
    # if ext in loaders:
    #     loader = loaders[ext](filepath)
    #     docs = loader.load()  # 数据转换
    #
    # func_name = f'extract_text_from_{ext}'
    # params = (f"{filepath}")
    # result = eval(f"{func_name}({', '.join(repr(param) for param in params))}")
    full_text = get_file_content(filepath)
    print(full_text)


    file_name=os.path.basename(filepath)


    file_name = os.path.basename(filepath)

    if emb_type == "openai":
        # 调用openai Embeddings  OPENAI_API_BASE
        # embeddings = OpenAIEmbeddings(openai_api_key=os.environ["OPENAI_API_KEY"])
        embeddings = OpenAIEmbeddings(openai_api_key="sk-FgmkVGYRirTVzJrjDMZ5Wi27ekHKq57xGHL2lZO6lTMuUAj3",
                                      openai_api_base="https://api.chatanywhere.tech/v1/")
        # embedding_model_name = 'shibing624/text2vec-bge-large-chinese'
    else:
        embeddings = HuggingFaceEmbeddings(model_name=embedding_model_name)

    # 文档切块目的是为了防止超出GPTAPI的token限制 RecursiveCharacterTextSplitter,CharacterTextSplitter,TokenTextSplitter,CodeTextSplitter
    # text_splitter = TokenTextSplitter(chunk_size=chunk_size, chunk_overlap=chunk_overlap)
    # text_splitter = RecursiveCharacterTextSplitter()
    # text_splitter = CharacterTextSplitter()

    # doc_texts = text_splitter.split_documents(docs)   #分隔文本


    text_splitter = SentenceSplitter(chunk_size=chunk_size, chunk_overlap=chunk_overlap)
    doc_texts = text_splitter.split_text(full_text)
    for doc in doc_texts:
        print(doc)
        doc.metadata["source"] = file_name
    # 向量化
    vectordb = Chroma.from_documents(doc_texts, embeddings, persist_directory=persist_directory)
    # 持久化
    vectordb.persist()
    # 执行到这里你会发现public目录下多了一些以parquest结尾的文件,这些文件就是chroma持久化本地的向量数据
    del embeddings


def update_vector(filepath, persist_directory, embedding_model_name, emb_type="openai", chunk_size=500,
                  chunk_overlap=20):
    delete_vector(filepath, persist_directory, embedding_model_name, emb_type)
    savevector(filepath, persist_directory, embedding_model_name, emb_type, chunk_size, chunk_overlap)


def delete_vector(filepath, persist_directory, embedding_model_name, emb_type="openai"):
    file_name = os.path.basename(filepath)
    # if emb_type == "openai":
    #     # 调用openai Embeddings
    #     embeddings = OpenAIEmbeddings(openai_api_key=os.environ["OPENAI_API_KEY"])
    #     # embedding_model_name = 'shibing624/text2vec-bge-large-chinese'
    # else:
    #     embeddings = HuggingFaceEmbeddings(model_name=embedding_model_name)
    persist_directory = os.path.abspath(persist_directory)
    # vectordb = Chroma(persist_directory=persist_directory, embedding_function=embeddings)
    vectordb = Chroma(persist_directory=persist_directory)
    docs = vectordb.get(where={"source": file_name})
    if len(docs["ids"]) > 0:
        vectordb.delete(ids=docs["ids"])
    # del embeddings


def getvectorkm_String(question, persist_directory, embedding_model_name, emb_type="openai"):
    if emb_type == "openai":
        # 调用openai Embeddings
        # embeddings = OpenAIEmbeddings(openai_api_key=os.environ["OPENAI_API_KEY"])
        embeddings = OpenAIEmbeddings(openai_api_key="sk-FgmkVGYRirTVzJrjDMZ5Wi27ekHKq57xGHL2lZO6lTMuUAj3",
                                      openai_api_base="https://api.chatanywhere.tech/v1/")
        # embedding_model_name = 'shibing624/text2vec-bge-large-chinese'
    else:
        embeddings = HuggingFaceEmbeddings(model_name=embedding_model_name)

    # 搜索
    # question = "对经国家、省、市等有关部门认定的企业技术中心及制造业创新中心，奖补政策是怎样的？"
    # 通过目录加载向量 这里的目录就是我们持久化的目录
    # persist_directory="C:\\dev\\ai-sns\\PyTalk\\pytalk\\vector_store\\vector"
    persist_directory = os.path.abspath(persist_directory)
    vectordb = Chroma(persist_directory=persist_directory, embedding_function=embeddings)
    # 向量搜索 根据你的问题进行本地向量搜索
    docs = vectordb.similarity_search_with_score(question, k=4)
    print(docs)
    # 将搜索到的信息重新转换为向量 (直接查到向量数据还不会😑😑）
    return docs


def get_file_content_tuple(file_path):
    # 判断文件扩展名
    _, file_extension = os.path.splitext(file_path)

    # 设定图片文件扩展名
    image_extensions = {'.jpg', '.jpeg', '.png', '.gif', '.bmp'}

    if file_extension.lower() in image_extensions:
        # 如果是图片文件，调用 image_to_base64 函数
        content = image_to_base64(file_path)
        return ("image", content, file_path)
    else:
        # 否则调用 get_file_content 函数
        content = get_file_content(file_path)
        return ("document", content, file_path)


def get_file_content(file_path):
    ext = os.path.splitext(file_path)[1].lower()  # 获取文件扩展名并转为小写
    exts = {
        '.js': 'txt',
        '.txt': 'txt',
        '.sql': 'txt',
        '.pdf': 'pdf',
        '.docx': 'docx',
        '.xls': 'xlsx',
        '.xlsx': 'xlsx',
        '.csv': 'csv',
        '.pptx': 'pptx',
        '.md': 'md',
        '.html': 'html',
        '.htm': 'html',
    }
    # func_name = f'extract_text_from_{ext}'
    # params = (f"{file_path}")
    # # result = eval(f"{func_name}({', '.join(repr(param) for param in params))}")
    # result = eval(f"{func_name}({repr(params[0])}")
    # print(result)
    # return result

    if ext in exts:
        func_name = f'extract_text_from_{exts[ext]}'
        result = eval(f"{func_name}({repr(f'{file_path}')})")
        return result

    raise ValueError(f"Unsupported file extension: {ext}")


# def get_file_content(file_path):
#     ext = os.path.splitext(file_path)[1].lower()  # 获取文件扩展名并转为小写
#
#     loaders = {
#         '.js': lambda path: TextLoader(path, encoding='utf8'),
#         '.txt': lambda path: TextLoader(path, encoding='utf8'),
#         '.sql': lambda path: TextLoader(path, encoding='utf8'),
#         '.pdf': PyPDFLoader,
#         '.docx': word_document.UnstructuredWordDocumentLoader,
#         '.xls': excel.UnstructuredExcelLoader,
#         '.xlsx': excel.UnstructuredExcelLoader,
#         '.csv': csv_loader.UnstructuredCSVLoader,
#         '.pptx': powerpoint.UnstructuredPowerPointLoader,
#         '.md': markdown.UnstructuredMarkdownLoader,
#         '.html': html.UnstructuredHTMLLoader,
#         '.htm': html.UnstructuredHTMLLoader,
#     }
#
#     if ext in loaders:
#         loader = loaders[ext](file_path)
#         docs = loader.load()  # 数据转换
#         return ''.join(doc.page_content for doc in docs)  # 连接所有的 page_content
#
#     raise ValueError(f"Unsupported file extension: {ext}")


class SentenceSplitter:
    def __init__(self, chunk_size: int = 250, chunk_overlap: int = 50):
        self.chunk_size = chunk_size
        self.chunk_overlap = chunk_overlap

    def split_text(self, text: str) -> List[str]:
        if self._is_has_chinese(text):
            return self._split_chinese_text(text)
        else:
            return self._split_english_text(text)

    def _split_chinese_text(self, text: str) -> List[str]:
        sentence_endings = {'\n', '。', '！', '？', '；', '…'}  # 句末标点符号
        chunks, current_chunk = [], ''
        for word in jieba.cut(text):
            if len(current_chunk) + len(word) > self.chunk_size:
                chunks.append(current_chunk.strip())
                current_chunk = word
            else:
                current_chunk += word
            if word[-1] in sentence_endings and len(current_chunk) > self.chunk_size - self.chunk_overlap:
                chunks.append(current_chunk.strip())
                current_chunk = ''
        if current_chunk:
            chunks.append(current_chunk.strip())
        if self.chunk_overlap > 0 and len(chunks) > 1:
            chunks = self._handle_overlap(chunks)
        return chunks

    def _split_english_text(self, text: str) -> List[str]:
        # 使用正则表达式按句子分割英文文本
        sentences = re.split(r'(?<=[.!?])\s+', text.replace('\n', ' '))
        chunks, current_chunk = [], ''
        for sentence in sentences:
            if len(current_chunk) + len(sentence) <= self.chunk_size or not current_chunk:
                current_chunk += (' ' if current_chunk else '') + sentence
            else:
                chunks.append(current_chunk)
                current_chunk = sentence
        if current_chunk:  # Add the last chunk
            chunks.append(current_chunk)

        if self.chunk_overlap > 0 and len(chunks) > 1:
            chunks = self._handle_overlap(chunks)

        return chunks

    def _is_has_chinese(self, text: str) -> bool:
        # check if contains chinese characters
        if any("\u4e00" <= ch <= "\u9fff" for ch in text):
            return True
        else:
            return False

    def _handle_overlap(self, chunks: List[str]) -> List[str]:
        # 处理块间重叠
        overlapped_chunks = []
        for i in range(len(chunks) - 1):
            chunk = chunks[i] + ' ' + chunks[i + 1][:self.chunk_overlap]
            overlapped_chunks.append(chunk.strip())
        overlapped_chunks.append(chunks[-1])
        return overlapped_chunks


@staticmethod
def extract_text_from_pdf(file_path: str):
    """Extract text content from a PDF file."""
    import PyPDF2
    contents = []
    with open(file_path, 'rb') as f:
        pdf_reader = PyPDF2.PdfReader(f)
        for page in pdf_reader.pages:
            page_text = page.extract_text().strip()
            raw_text = [text.strip() for text in page_text.splitlines() if text.strip()]
            new_text = ''
            for text in raw_text:
                new_text += text
                if text[-1] in ['.', '!', '?', '。', '！', '？', '…', ';', '；', ':', '：', '”', '’', '）', '】', '》', '」',
                                '』', '〕', '〉', '》', '〗', '〞', '〟', '»', '"', "'", ')', ']', '}']:
                    contents.append(new_text)
                    new_text = ''
            if new_text:
                contents.append(new_text)
    data = ','.join(contents)
    return data


@staticmethod
def extract_text_from_txt(file_path: str):
    """Extract text content from a TXT file."""
    with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
        contents = [text.strip() for text in f.readlines() if text.strip()]
    data = ','.join(contents)
    return data


@staticmethod
def extract_text_from_docx(file_path: str):
    """Extract text content from a DOCX file."""
    import docx
    document = docx.Document(file_path)
    contents = [paragraph.text.strip() for paragraph in document.paragraphs if paragraph.text.strip()]
    data = ','.join(contents)
    return data


@staticmethod
def extract_text_from_markdown(file_path: str):
    """Extract text content from a Markdown file."""
    import markdown
    from bs4 import BeautifulSoup
    with open(file_path, 'r', encoding='utf-8') as f:
        markdown_text = f.read()
    html = markdown.markdown(markdown_text)
    soup = BeautifulSoup(html, 'html.parser')
    contents = [text.strip() for text in soup.get_text().splitlines() if text.strip()]
    data = ','.join(contents)
    return data


# 打开CSV文件
@staticmethod
def extract_text_from_csv(file_path: str):
    import csv
    contents = []
    with open(file_path, newline='', encoding='utf-8') as csvfile:
        # 创建一个csv阅读器
        csvreader = csv.reader(csvfile)
        # 遍历CSV文件中的每一行
        for row in csvreader:
            # 将每行的元素合并为一个字符串，元素之间用逗号分隔
            text = ','.join(row)
            # 打印或返回文本
            contents.append(text)
    data = ','.join(contents)
    return data


@staticmethod
def extract_text_from_xlsx(file_path: str):
    """Extract data from an Excel file."""
    # 使用pandas读取Excel文件，默认读取第一个工作表
    contents = []
    import pandas as pd
    df = pd.read_excel(file_path, engine='openpyxl')
    # 将DataFrame转换为列表，其中每个元素是一行数据
    res = df.values.tolist()
    data = df.to_string(index=False)
    # data = '\n'.join([','.join(map(str, row)) for row in res])
    return data


# 调用函数并传入PPTX文件路径
@staticmethod
def extract_text_from_pptx(file_path: str):
    from pptx import Presentation
    """Extract text content from a PPTX file."""
    # 加载PPTX文件 python-pptx
    presentation = Presentation(file_path)
    # 初始化一个列表来存储文本内容
    contents = []
    # 遍历所有幻灯片
    for slide in presentation.slides:
        # 遍历幻灯片中的所有形状
        for shape in slide.shapes:
            # 如果形状包含文本框
            if hasattr(shape, "text"):
                # 获取文本框中的文本，并去除首尾空格
                text = shape.text.strip()
                # 如果文本不为空，则添加到列表中
                if text:
                    contents.append(text)

    data = ','.join(contents)
    return data


@staticmethod
def extract_text_from_html(file_path: str):
    """Extract text content from an HTML file."""
    # 打开HTML文件
    from bs4 import BeautifulSoup
    with open(file_path, 'r', encoding='utf-8', errors='ignore') as file:
        # 读取文件内容
        html_content = file.read()

    # 使用BeautifulSoup解析HTML内容
    soup = BeautifulSoup(html_content, 'lxml')

    # 提取并返回文本内容
    # .get_text() 方法可以提取标签内的文本，separator参数用于指定文本之间的分隔符
    return soup.get_text(separator=' ', strip=True)
